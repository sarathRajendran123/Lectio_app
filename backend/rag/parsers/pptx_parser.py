"""
LECTIO — PPTX Parser (python-pptx)
Extracts slide titles, body text, and speaker notes.

Design decisions:
  - Slide ORDER is preserved — critical for Content-Delivery alignment auditing.
    A topic on slide 32 that prerequisites slide 40 content is a sequencing bug.
  - Speaker notes extracted separately (block_type="note") — often contain
    richer explanations than the slide bullets.
  - Each slide title becomes a section_title for all subsequent blocks from
    that slide — enabling week/topic inference downstream.
"""

import logging
from pathlib import Path
from typing import List, Optional

from rag.parsers.base_parser import BaseParser, ParsedBlock, ParsedDocument

logger = logging.getLogger(__name__)


class PptxParser(BaseParser):

    def parse(self, file_path: str, artifact_type: str = "slides") -> ParsedDocument:
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx not installed. Run: pip install python-pptx")

        path = Path(file_path)
        prs  = Presentation(str(path))

        blocks: List[ParsedBlock] = []
        warnings: List[str] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_title: Optional[str] = None

            # ── Extract title ─────────────────────────────────────────────────
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title_text = slide.shapes.title.text_frame.text.strip()
                if title_text:
                    slide_title = title_text
                    blocks.append(ParsedBlock(
                        text=self._safe_text(title_text),
                        block_type="title",
                        slide_number=slide_num,
                        section_title=slide_title,
                        heading_level=1,
                        is_bold=True,
                    ))

            # ── Extract body text from all shapes ─────────────────────────────
            for shape in slide.shapes:
                # Skip the title shape (already handled)
                if shape == slide.shapes.title:
                    continue

                if not shape.has_text_frame:
                    continue

                shape_text_parts = []
                for para in shape.text_frame.paragraphs:
                    para_text = para.text.strip()
                    if para_text:
                        shape_text_parts.append(para_text)

                if not shape_text_parts:
                    continue

                combined = "\n".join(shape_text_parts)
                combined = self._safe_text(combined)
                if not combined:
                    continue

                blocks.append(ParsedBlock(
                    text=combined,
                    block_type="paragraph",
                    slide_number=slide_num,
                    section_title=slide_title,
                    metadata={"shape_name": shape.name},
                ))

            # ── Extract speaker notes ─────────────────────────────────────────
            try:
                notes_frame = slide.notes_slide.notes_text_frame if slide.has_notes_slide else None
                if notes_frame:
                    notes_text = notes_frame.text.strip()
                    if notes_text and len(notes_text) > 10:
                        blocks.append(ParsedBlock(
                            text=self._safe_text(notes_text),
                            block_type="note",
                            slide_number=slide_num,
                            section_title=slide_title,
                            metadata={"is_speaker_note": True},
                        ))
            except Exception as e:
                warnings.append(f"Slide {slide_num}: could not extract speaker notes ({e})")

        all_text   = " ".join(b.text for b in blocks)
        word_count = self._count_words(all_text)

        return ParsedDocument(
            source_path=str(path),
            file_type="pptx",
            artifact_type=artifact_type,
            blocks=blocks,
            total_slides=len(prs.slides),
            title=path.stem,
            word_count=word_count,
            parse_warnings=warnings,
        )
